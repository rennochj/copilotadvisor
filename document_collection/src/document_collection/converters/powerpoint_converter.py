"""PowerPoint to Markdown converter implementation."""

import zipfile
from pathlib import Path
from typing import Any

from document_collection.core.interfaces import DocumentConverter


class PowerPointConverter(DocumentConverter):
    """Convert PowerPoint documents to Markdown format with image extraction."""

    def can_convert(self, file_path: Path) -> bool:
        """Check if this converter can handle PowerPoint files."""
        return str(file_path).lower().endswith(".pptx")

    def get_supported_formats(self) -> list[str]:
        """Return list of supported formats."""
        return ["pptx"]

    async def convert(self, input_path: Path, output_path: Path, **kwargs: Any) -> Path:
        """Convert PowerPoint document to Markdown using python-pptx with image extraction."""
        try:
            from pptx import Presentation

            # Read PowerPoint presentation
            prs = Presentation(str(input_path))

            # Create images directory
            images_dir = output_path.parent / "images"
            images_dir.mkdir(parents=True, exist_ok=True)

            # Extract embedded images from the presentation's media folder
            image_counter = 1
            image_mapping = {}  # Map original image names to new filenames

            try:
                # Access the presentation as a zip file to extract images
                with zipfile.ZipFile(str(input_path), 'r') as pptx_zip:
                    # Look for images in the media directory
                    for file_info in pptx_zip.filelist:
                        if file_info.filename.startswith('ppt/media/'):
                            # Extract image
                            image_data = pptx_zip.read(file_info.filename)

                            # Determine file extension
                            original_name = Path(file_info.filename).name
                            file_ext = Path(file_info.filename).suffix.lower()
                            if not file_ext:
                                file_ext = '.png'  # Default

                            # Create new filename
                            new_filename = f"{input_path.stem}_image_{image_counter:03d}{file_ext}"
                            image_path = images_dir / new_filename

                            # Save image
                            with open(image_path, 'wb') as img_file:
                                img_file.write(image_data)

                            # Map original to new filename
                            image_mapping[original_name] = new_filename
                            image_counter += 1

            except Exception:
                # Continue if image extraction fails
                pass

            # Convert presentation content to markdown
            content_parts = [f"# {input_path.stem}\n"]
            content_parts.append("Converted from PowerPoint presentation\n")

            if image_counter > 1:
                content_parts.append(f"*This presentation contains {image_counter - 1} extracted images stored in the `images/` directory.*\n")

            # Process slides
            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts.append(f"\n## Slide {slide_num}\n")

                # Extract text from shapes
                for shape in slide.shapes:
                    try:
                        # Try to get text from text_frame first
                        text_frame = getattr(shape, 'text_frame', None)
                        if text_frame:
                            text = getattr(text_frame, 'text', '').strip()
                            if text:
                                # Add text with some basic formatting
                                if len(text) < 100 and '\n' not in text:
                                    # Likely a title or header
                                    content_parts.append(f"### {text}\n")
                                else:
                                    # Body text
                                    content_parts.append(f"{text}\n")
                        else:
                            # Fallback to direct text attribute
                            text = getattr(shape, 'text', '').strip()
                            if text:
                                # Add text with some basic formatting
                                if len(text) < 100 and '\n' not in text:
                                    # Likely a title or header
                                    content_parts.append(f"### {text}\n")
                                else:
                                    # Body text
                                    content_parts.append(f"{text}\n")
                    except Exception:
                        # Skip shapes that don't have text
                        pass

                    # Check for tables in shape
                    try:
                        table = getattr(shape, 'table', None)
                        if table:
                            content_parts.append("\n")
                            for i, row in enumerate(table.rows):
                                row_cells = [getattr(cell, 'text', '').strip() for cell in row.cells]
                                content_parts.append("| " + " | ".join(row_cells) + " |\n")

                                # Add header separator for first row
                                if i == 0:
                                    content_parts.append("| " + " | ".join(["---"] * len(row_cells)) + " |\n")
                            content_parts.append("\n")
                    except Exception:
                        # Skip if table extraction fails
                        pass

                # Add notes if present
                notes_slide = getattr(slide, 'notes_slide', None)
                if notes_slide:
                    notes_text = ""
                    try:
                        for shape in notes_slide.shapes:
                            try:
                                text_frame = getattr(shape, 'text_frame', None)
                                if text_frame:
                                    notes_text += getattr(text_frame, 'text', '').strip() + " "
                                else:
                                    notes_text += getattr(shape, 'text', '').strip() + " "
                            except Exception:
                                continue
                        if notes_text.strip():
                            content_parts.append(f"\n**Speaker Notes:** {notes_text.strip()}\n")
                    except Exception:
                        # Skip notes if extraction fails
                        pass

            # Add image references for any images found
            if image_mapping:
                content_parts.append("\n## Extracted Images\n\n")
                for i, (_original, new_filename) in enumerate(image_mapping.items(), 1):
                    content_parts.append(f"![Image {i}](images/{new_filename})\n\n")

            # Write markdown content
            markdown_content = "\n".join(content_parts)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)

            return output_path

        except ImportError:
            # Fallback if python-pptx is not available
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(f"# {input_path.stem}\n\nPowerPoint conversion requires python-pptx library.\n")
            return output_path
        except Exception as e:
            # Create error file
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(f"# {input_path.stem}\n\nError converting PowerPoint document: {str(e)}\n")
            return output_path
